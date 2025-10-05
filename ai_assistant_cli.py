# AI Assistant CLI - Final Project for Beginners
# This tool demonstrates a menu-driven Python CLI using the Gemini API (via HTTP requests)
# for Q&A, file processing, and command management.

import requests
import json
import os
import time
import zipfile
import tarfile
import base64
from io import BytesIO
from uuid import uuid4 # Used for generating unique IDs for recipes

# --- Library Installation and Imports ---

# In environments like Google Colab, we must explicitly install external libraries.
try:
    print("Initializing environment and installing necessary libraries...")
    # Libraries for PDF, DOCX, XLSX, PPTX, and Image handling are installed here.
    # The output is suppressed with '> /dev/null 2>&1' for a cleaner terminal experience.
    os.system("pip install PyMuPDF python-docx openpyxl python-pptx Pillow > /dev/null 2>&1")
    print("All file processing libraries installed/verified.")
except Exception as e:
    print(f"Could not automatically install libraries: {e}")

# Import third-party libraries after the installation attempt
try:
    import fitz # PyMuPDF for PDF
    import docx # python-docx
    from openpyxl import load_workbook
    from pptx import Presentation # python-pptx
    from PIL import Image
except ImportError as e:
    # This warning helps the user debug if installation failed.
    print(f"FATAL WARNING: One or more required libraries failed to import: {e}. Summarization may fail.")


# --- Configuration ---
# FIX: Key is now loaded from the environment variable GEMINI_API_KEY.
# This requires running the %env command in a Colab cell before the script.
RAW_API_KEY = os.environ.get("GEMINI_API_KEY", "YOUR_KEY_HERE")
API_KEY = RAW_API_KEY.strip("'\"") # Strip any extra quotes the environment might add
API_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-preview-05-20:generateContent"
GROUNDED_MODEL_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-preview-05-20:generateContent"
COMMAND_RECIPES_FILE = 'command_recipes.json' # File to store command recipes


# --- Core API Interaction ---

def _call_gemini_api(payload, is_grounded=False, image_data_base64=None):
    """Handles the request to the Gemini API with support for image data and exponential backoff."""

    if API_KEY == "YOUR_KEY_HERE":
        return "Error: API Key is not set. Please set the GEMINI_API_KEY environment variable.", []

    url = GROUNDED_MODEL_URL if is_grounded else API_URL

    headers = {
        'Content-Type': 'application/json',
        'X-Goog-Api-Key': API_KEY,
    }

    # 1. Prepare contents array (Handles multimodal input)
    contents_parts = payload.get('contents', [{}])[0].get('parts', [])

    if image_data_base64:
        # Insert image data at the beginning of the parts list for vision models
        image_part = {
            "inlineData": {
                # We save images as JPEG for consistency
                "mimeType": "image/jpeg", 
                "data": image_data_base64
            }
        }
        contents_parts.insert(0, image_part)
        payload['contents'] = [{"parts": contents_parts}]


    max_retries = 3
    for attempt in range(max_retries):
        try:
            print("... Sending request to AI model...")
            # Adding verify=False to bypass potential Colab/local SSL issues
            requests.packages.urllib3.disable_warnings(requests.packages.urllib3.exceptions.InsecureRequestWarning)
            response = requests.post(url, headers=headers, data=json.dumps(payload), verify=False)
            response.raise_for_status() # Raises HTTPError for bad responses (4xx or 5xx)

            result = response.json()
            candidate = result.get('candidates', [{}])[0]

            if candidate and candidate.get('content') and candidate['content'].get('parts'):
                text = candidate['content']['parts'][0].get('text', 'No response text found.')

                # Extract grounding sources (for web search)
                sources = []
                grounding_metadata = candidate.get('groundingMetadata', {})
                if grounding_metadata and grounding_metadata.get('groundingAttributions'):
                    sources = [
                        f"- [{attr.get('web', {}).get('title', 'Source')}]({attr.get('web', {}).get('uri', '#')})"
                        for attr in grounding_metadata['groundingAttributions']
                    ]
                return text, sources

            return "Error: AI response candidate was empty or malformed.", []

        except requests.exceptions.HTTPError as e:
            if response.status_code == 429 and attempt < max_retries - 1:
                # Handle rate limiting (429) with exponential backoff
                wait_time = 2 ** attempt
                print(f"Rate limit hit (429). Retrying in {wait_time} seconds...")
                time.sleep(wait_time)
            else:
                return f"HTTP Error: {e}", []
        except requests.exceptions.RequestException as e:
            return f"Network Error: {e}", []
        except json.JSONDecodeError:
            return "Error: Failed to decode JSON response from API.", []

    return "Error: Failed to get a response after multiple retries.", []


# --- Document & Archive Extraction Helpers (omitted for brevity) ---

def _extract_text_from_proprietary_docs(file_path: str) -> str:
    """Extracts text from DOCX, XLSX, and PPTX files."""
    ext = os.path.splitext(file_path)[1].lower()
    content = []

    try:
        if ext == '.docx':
            doc = docx.Document(file_path)
            content.append("\n".join([p.text for p in doc.paragraphs]))

        elif ext == '.xlsx':
            workbook = load_workbook(file_path)
            for sheet_name in workbook.sheetnames:
                content.append(f"\n--- Spreadsheet: {sheet_name} ---\n")
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows():
                    content.append(", ".join([str(cell.value) if cell.value is not None else "" for cell in row]))
            content.append("\n")

        elif ext == '.pptx':
            presentation = Presentation(file_path)
            for slide_number, slide in enumerate(presentation.slides):
                content.append(f"\n--- Slide {slide_number + 1} (Title: {slide.shapes.title.text if slide.has_title else 'N/A'}) ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content.append(shape.text)
            content.append("\n")

        return "\n".join(content)

    except Exception as e:
        return f"Error: Failed to process proprietary file ({ext}): {e}"


def _extract_text_from_archive(file_path: str) -> str:
    """Extracts the file listing and structure from .zip, .tar, .gz archives."""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == '.zip':
            with zipfile.ZipFile(file_path, 'r') as archive:
                file_list = "\n".join(archive.namelist())
                return f"Archive Type: ZIP\nContents:\n{file_list}"

        elif ext in ('.tar', '.gz', '.tgz', '.bz2', '.tar.gz'):
            mode = 'r' if ext == '.tar' else 'r:*'
            with tarfile.open(file_path, mode) as archive:
                file_list = "\n".join(archive.getnames())
                return f"Archive Type: TAR\nContents:\n{file_list}"

        return f"Error: Unsupported archive format: {ext}"

    except FileNotFoundError:
        return f"Error: Archive file not found at path: {file_path}"
    except Exception as e:
        return f"Error processing archive file: {e}"


def _get_image_base64(file_path: str) -> tuple:
    """Converts image file to Base64 string for Gemini Vision API."""
    try:
        # Resize image to save bandwidth and stay within typical size limits
        MAX_SIZE = (1024, 1024)

        img = Image.open(file_path)
        img.thumbnail(MAX_SIZE) # Resize in place

        # Convert image to bytes in JPEG format (universal format for API)
        buffer = BytesIO()
        img.save(buffer, format="JPEG")

        # Encode bytes to Base64 string
        img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        return img_base64, None # Return image data and no error

    except FileNotFoundError:
        return None, "Error: Image file not found."
    except Exception as e:
        return None, f"Error converting image to Base64: {e}"


def _extract_text_from_plain_and_pdf(file_path: str) -> str:
    """Handles standard text and PDF files."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.pdf':
        try:
            with fitz.open(file_path) as pdf_document:
                text_content = "".join(page.get_text() for page in pdf_document)
                if not text_content.strip():
                     return f"Warning: PDF file read, but no readable text extracted."
                return text_content
        except NameError:
            return "Error: PyMuPDF (fitz) library not imported. PDF summarization failed."
        except Exception as e:
            return f"Error processing PDF: {e}"

    # Handle standard text files (all other types)
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            # Fallback to Latin-1 for files with non-standard characters
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            return f"Error: Could not read file using UTF-8 or Latin-1 encoding: {e}"
    except Exception as e:
        return f"An unknown error occurred while reading the file: {e}"

def summarize_file():
    """Reads a local file (many types) and asks the AI to summarize its contents."""
    print("\n--- 4. Summarize Local File (Multi-format) ---") # Corrected menu number
    filepath = input("Enter the path to the file you want to analyze (or type 'back'):\n> ").strip()
    if filepath.lower() == 'back':
        return

    if not os.path.exists(filepath):
        print(f"Error: File not found at path: {filepath}")
        return

    ext = os.path.splitext(filepath)[1].lower()
    file_content = ""
    image_base64 = None

    if ext in ('.txt', '.md', '.log', '.py', '.js', '.json', '.yaml', '.yml', '.csv', '.pdf'):
        file_content = _extract_text_from_plain_and_pdf(filepath)
    elif ext in ('.zip', '.tar', '.gz', '.tgz', '.bz2', '.tar.gz'):
        file_content = _extract_text_from_archive(filepath)
    elif ext in ('.docx', '.xlsx', '.pptx'):
        file_content = _extract_text_from_proprietary_docs(filepath)
    elif ext in ('.jpg', '.jpeg', '.png', '.webp'):
        image_base64, error = _get_image_base64(filepath)
        if error:
            print(error)
            return
        file_content = "Describe and summarize this image content in a concise, bulleted list."
    else:
        print(f"Error: Unsupported file type for direct analysis: {ext}. Returning to menu.")
        return

    if file_content.startswith("Error") or file_content.startswith("Warning"):
        print(file_content)
        return

    if not image_base64:
        print(f"Successfully extracted {len(file_content)} characters from '{filepath}'.")
        # For text files, the prompt includes the content
        prompt = f"Please provide a concise, bulleted summary of the following document content:\n\n---\n{file_content}\n---"
    else:
        # For images, the prompt is simple and the image is passed via Base64
        prompt = file_content 

    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
    }

    response_text, _ = _call_gemini_api(payload, image_data_base64=image_base64)

    print("\n" + "="*50)
    print(f"📄 Summary of {os.path.basename(filepath)}:")
    print(response_text)
    print("="*50 + "\n")


# --- Recipe Book Persistence and AI Helpers ---

def _load_recipes():
    """Loads command recipes from the local JSON file."""
    if not os.path.exists(COMMAND_RECIPES_FILE):
        return []
    try:
        with open(COMMAND_RECIPES_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"Error loading recipes: {e}. Starting with an empty list.")
        return []

def _save_recipes(recipes):
    """Saves command recipes to the local JSON file."""
    try:
        with open(COMMAND_RECIPES_FILE, 'w', encoding='utf-8') as f:
            json.dump(recipes, f, indent=4)
        print(f"\n✅ Recipes saved to {COMMAND_RECIPES_FILE}.")
    except Exception as e:
        print(f"Error saving recipes: {e}")

def _get_ai_suggested_tags(command: str) -> list:
    """Uses Gemini to suggest tags for a given command string."""
    prompt = (
        "Analyze the following Linux/CLI command. Provide exactly 5 relevant tags, separated by commas. "
        "Do NOT include any explanation, headers, or extra text, only the tags."
        f"Command: {command}"
    )

    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
    }

    print("... AI analyzing command to suggest tags...")
    response_text, _ = _call_gemini_api(payload)

    if response_text.startswith("Error"):
        print(f"Warning: AI tagging failed due to API error. Using no suggested tags.")
        return []

    cleaned_tags = [
        tag.strip().lower() for tag in response_text.replace('\n', '').split(',')
        if tag.strip()
    ]
    return cleaned_tags


# The function that generates command and name is now simplified to just return the command string
def _get_ai_generated_command(description: str) -> str:
    """Uses Gemini to generate ONLY the CLI command string from a description."""
    prompt = (
        "You are an expert Linux/macOS command line utility expert. Your sole output must be a single, complete, "
        "executable command line string that fulfills the user's request. DO NOT include any explanatory text, "
        "markdown formatting, or extra characters. Only provide the command itself."
        f"Request: {description}"
    )

    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
    }

    print("... AI generating command...")
    response_text, _ = _call_gemini_api(payload)
    
    if response_text.startswith("Error"):
        print(f"Error: AI command generation failed: {response_text}")
        return None
    
    # Clean up any residual markdown or explanation the AI might have accidentally added
    generated_command = response_text.strip().split('\n')[0].replace('`', '').replace('bash', '').strip()
    
    return generated_command


# --- AI Assistant Functions ---

def chat_with_grounding():
    """General Q&A using the Gemini model with Google Search grounding."""
    print("\n--- 1. General AI Chat (Web Search) ---")
    query = input("Ask your question (or type 'back'):\n> ")
    if query.lower() == 'back':
        return

    payload = {
        "contents": [{"parts": [{"text": query}]}],
        "tools": [{"google_search": {}}] # Enable Google Search for grounding
    }
    
    response_text, sources = _call_gemini_api(payload, is_grounded=True)
    
    print("\n" + "="*50)
    print("🤖 AI Response:")
    print(response_text)
    
    if sources:
        print("\n🌐 Sources (Grounding):")
        for source in sources:
            print(source)
    print("="*50 + "\n")


def code_explainer():
    """Generates an explanation for a user-provided code snippet."""
    print("\n--- 2. Code Explainer ---")
    print("Paste your code snippet below. Type 'END' on a new line when finished.")
    
    code_lines = []
    while True:
        line = input()
        if line.strip().upper() == 'END':
            break
        code_lines.append(line)
        
    code_snippet = "\n".join(code_lines)
    
    if not code_snippet.strip():
        print("No code provided. Returning to main menu.")
        return
        
    prompt = f"Explain the following code snippet thoroughly, focusing on its purpose, inputs, and outputs. Use clear, simple language:\n\n\n{code_snippet}\n"

    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
    }

    response_text, _ = _call_gemini_api(payload)
    
    print("\n" + "="*50)
    print("🧠 Code Analysis:")
    print(response_text)
    print("="*50 + "\n")


def code_generator():
    """Generates code based on a natural language request."""
    print("\n--- 3. Code Generator ---")
    request = input("Describe the code you need (e.g., 'Python script for a merge sort algorithm'):\n> ").strip()
    if not request:
        return

    # System instruction to encourage clean, runnable code wrapped in a Markdown block
    system_prompt = "You are an expert software developer. Generate a complete, runnable code solution for the user's request. Always wrap the code in a single Markdown code block (language\\ncode\\n)."
    
    prompt = f"Generate code for the following request: {request}"

    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "systemInstruction": {"parts": [{"text": system_prompt}]}
    }

    response_text, _ = _call_gemini_api(payload)
    
    print("\n" + "="*50)
    print("💡 Generated Code:")
    print(response_text)
    print("="*50 + "\n")


# --- Recipe Book Functions ---

def suggest_new_command_recipe():
    """Uses AI to generate a command based on user description and saves it."""
    print("\n--- 5. Generate New Recipe (AI) ---") # Corrected menu number
    description = input("Describe the task you want a CLI command for (e.g., 'how to delete all old docker images'):\n> ").strip()

    if not description:
        print("Description cannot be empty. Returning to menu.")
        return

    # Use the simplified function to get JUST the command string
    command = _get_ai_generated_command(description)

    if not command:
        print("Failed to generate a valid command. Please try a different description.")
        return

    # Use the description as the base name for the recipe
    name = f"AI Generated: {description[:50]}..." if len(description) > 50 else description

    print("\n" + "="*50)
    print("🤖 AI Suggestion:")
    print(f"Name: {name}")
    print(f"Command: {command}")
    print("="*50)

    confirm = input("Do you want to save this command to your Recipe Book? (y/n): ").strip().lower()

    if confirm == 'y':
        suggested_tags = _get_ai_suggested_tags(command)
        
        print("\n--- Tagging ---")
        if suggested_tags:
            print(f"💡 AI Suggested Tags: {', '.join(suggested_tags)}")
            tags_input = input("Enter additional tags (or press Enter to use AI's suggestions):\n> ").strip()
        else:
            tags_input = input("Enter comma-separated tags:\n> ").strip()

        user_tags = [t.strip().lower() for t in tags_input.split(',') if t.strip()]
        final_tags = list(set(suggested_tags + user_tags))

        recipes = _load_recipes()
        
        new_recipe = {
            'id': str(uuid4()),
            'name': name,
            'command': command,
            'tags': final_tags,
            'timestamp': time.strftime('%Y-%m-%d %H:%M:%S')
        }

        recipes.append(new_recipe)
        _save_recipes(recipes)
        print(f"Successfully saved AI-generated recipe '{name}'.")
    else:
        print("Recipe discarded.")


def add_known_command_recipe():
    """Adds a new complex command recipe (manual input) to the local vault, with AI suggested tags."""
    print("\n--- 6. Add Known Recipe (Manual) ---") # Corrected menu number
    name = input("Enter a unique name for this recipe (e.g., 'Grep Python Files'):\n> ").strip()
    if not name:
        print("Recipe name cannot be empty.")
        return

    command = input("Paste the full command string (e.g., 'grep -r --include=*.py \"import\" .'):\n> ").strip()
    if not command:
        print("Command cannot be empty.")
        return

    recipes = _load_recipes()

    if any(r['name'].lower() == name.lower() for r in recipes):
        print(f"Error: Recipe named '{name}' already exists. Use a unique name.")
        return

    suggested_tags = _get_ai_suggested_tags(command)
    
    print("\n--- Tagging ---")
    if suggested_tags:
        print(f"💡 AI Suggested Tags: {', '.join(suggested_tags)}")
        tags_input = input("Enter additional comma-separated tags (or press Enter to use AI's suggestions):\n> ").strip()
    else:
        tags_input = input("Enter comma-separated tags (AI suggestion failed):\n> ").strip()

    user_tags = [t.strip().lower() for t in tags_input.split(',') if t.strip()]
    final_tags = list(set(suggested_tags + user_tags))


    new_recipe = {
        'id': str(uuid4()),
        'name': name,
        'command': command,
        'tags': final_tags,
        'timestamp': time.strftime('%Y-%m-%d %H:%M:%S')
    }

    recipes.append(new_recipe)
    _save_recipes(recipes)
    print(f"Successfully added recipe '{name}' with {len(final_tags)} tags.")


def search_and_copy_recipe():
    """Searches for a command recipe and displays it."""
    print("\n--- 7. Search & View Recipe ---") # Corrected menu number
    recipes = _load_recipes()
    if not recipes:
        print("\nNo command recipes found. Please add one first (Option 5 or 6).")
        return

    query = input("Enter a keyword or tag to search:\n> ").strip().lower()

    results = [
        r for r in recipes
        if query in r['name'].lower() or
           any(query in tag for tag in r['tags'])
    ]

    if not results:
        print(f"No recipes found matching '{query}'.")
        return

    print(f"\nFound {len(results)} recipes:")
    for i, r in enumerate(results):
        print(f"{i+1}. Name: {r['name']}")
        print(f"  Tags: {', '.join(r['tags'])}")
        print(f"  Command: {r['command'][:70]}...") # Display truncated command

    try:
        selection = int(input("Enter the number of the recipe you want to view/copy:\n> ")) - 1
        if 0 <= selection < len(results):
            selected_recipe = results[selection]

            print("\n" + "#"*60)
            print(f"Recipe: {selected_recipe['name']}")
            print(f"Tags: {', '.join(selected_recipe['tags'])}")
            print("\nFull Command:")
            print("--------------------------------------------------")
            print(selected_recipe['command'])
            print("--------------------------------------------------")
            print("\n📌 NOTE: Please copy the command above and paste it into your shell.")
            print("#"*60 + "\n")
        else:
            print("Invalid selection number.")
    except ValueError:
        print("Invalid input. Please enter a number.")


def explain_command_recipe():
    """Searches for a command recipe and asks the AI to explain it."""
    recipes = _load_recipes()
    if not recipes:
        print("\nNo command recipes found. Please add one first (Option 5 or 6).")
        return

    print("\n--- 8. Explain Recipe (AI) ---") # Corrected menu number
    query = input("Enter keyword or tag to find the recipe you want explained:\n> ").strip().lower()

    results = [
        r for r in recipes
        if query in r['name'].lower() or
           any(query in tag for tag in r['tags'])
    ]

    if not results:
        print(f"No recipes found matching '{query}'.")
        return

    print(f"\nFound {len(results)} recipes:")
    for i, r in enumerate(results):
        print(f"{i+1}. Name: {r['name']} | Command: {r['command'][:50]}...")

    try:
        selection = int(input("Enter the number of the recipe to EXPLAIN (or type 0 to cancel):\n> ")) - 1
        if selection < 0 or selection >= len(results):
            print("Cancelled or Invalid selection.")
            return

        selected_recipe = results[selection]
        command_to_explain = selected_recipe['command']

        prompt = f"Act as a Linux/DevOps tutor. Explain the following command line recipe step-by-step. Focus on what each flag and argument does, and provide a clear, easy-to-understand purpose for the entire command. The command is:\n\n`{command_to_explain}`"

        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
        }

        response_text, _ = _call_gemini_api(payload)

        print("\n" + "="*50)
        print(f"🧠 AI Explanation for: {selected_recipe['name']}")
        print(f"Command: {command_to_explain}")
        print("--------------------------------------------------")
        print(response_text)
        print("="*50 + "\n")

    except ValueError:
        print("Invalid input. Please enter a number.")


# --- Main Menu and Execution ---

def main_menu():
    """Displays the main menu and handles user input."""
    if API_KEY == "YOUR_KEY_HERE":
        print("\nFATAL ERROR: Gemini API key is not set.")
        print("Please replace 'YOUR_KEY_HERE' in the code with your actual API key.")
        return

    while True:
        print("\n" + "╭───────────────────────╮")
        print("│ 🤖 AI Assistant & Recipes │")
        print("╰───────────────────────╯")
        print("--- AI Analysis & Generation ---")
        print("1. 🗣 General Q&A (Web Search)")
        print("2. 💻 Explain Code Snippet")
        print("3. 💡 Code Generator (NEW!)")
        print("4. 📝 Summarize Local File (Multi-format)")
        print("--- Command Recipe Vault ---")
        print("5. 🛠 Generate New Recipe (AI)")
        print("6. ➕ Add Known Recipe (Manual)")
        print("7. 🔍 Search & View Recipe")
        print("8. 🧠 Explain Recipe (AI)")
        print("9. 🚪 Exit")

        choice = input("Enter your choice (1-9): ")

        if choice == '1':
            chat_with_grounding()
        elif choice == '2':
            code_explainer()
        elif choice == '3':
            code_generator()
        elif choice == '4':
            summarize_file()
        elif choice == '5':
            suggest_new_command_recipe()
        elif choice == '6':
            add_known_command_recipe()
        elif choice == '7':
            search_and_copy_recipe()
        elif choice == '8':
            explain_command_recipe()
        elif choice == '9':
            print("Thank you for using the AI Assistant CLI. Goodbye!")
            break
        else:
            print("Invalid choice. Please enter a number between 1 and 9.")

if __name__ == "__main__":
    main_menu()